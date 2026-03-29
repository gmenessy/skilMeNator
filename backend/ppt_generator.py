import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Komm.ONE Colors
COLOR_MIDNIGHT = RGBColor(0, 58, 64) # #003A40
COLOR_LAGOON = RGBColor(0, 178, 169) # #00B2A9
COLOR_AMARILLO = RGBColor(241, 196, 0) # #F1C400
COLOR_WHITE = RGBColor(255, 255, 255)

def apply_komm_one_theme(prs):
    """
    Applies custom styling to the presentation based on CI tokens
    Note: Real custom master slides are best done by loading a base template.
    We'll do our best with code formatting here.
    """
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = 'DM Sans' # Or 'Arial' as fallback
                    if slide == prs.slides[0] and shape == slide.shapes[0]:
                        # Title slide title
                        run.font.color.rgb = COLOR_MIDNIGHT
                    elif shape == slide.shapes[0]:
                        # Normal slide title
                        run.font.color.rgb = COLOR_LAGOON

def create_presentation(slide_data, filename="presentation.pptx", output_dir="generated_files"):
    """
    Generates a PowerPoint presentation from structured JSON data.
    Expected data structure:
    {
      "title": "Presentation Title",
      "subtitle": "Presentation Subtitle",
      "slides": [
        {
          "title": "Slide 1 Title",
          "content": ["Point 1", "Point 2", "Point 3"]
        }
      ]
    }
    """
    prs = Presentation()

    # Create Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = slide_data.get("title", "Generated Presentation")
    subtitle.text = slide_data.get("subtitle", "Powered by ChatOrchestrator")

    # Create Content Slides
    for slide_info in slide_data.get("slides", []):
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes

        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = slide_info.get("title", "Slide")

        tf = body_shape.text_frame
        tf.clear() # Clear default prompt text

        content = slide_info.get("content", [])

        # Handle string content instead of array
        if isinstance(content, str):
            p = tf.paragraphs[0]
            p.text = content
            p.font.size = Pt(18)
        else:
            for idx, point in enumerate(content):
                if idx == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()

                # Check for sub-bullets
                if isinstance(point, dict) and 'text' in point:
                    p.text = point['text']
                    p.level = point.get('level', 0)
                else:
                    p.text = str(point)
                    p.level = 0

                p.font.size = Pt(18)

    # Attempt to style via code (best effort without a template)
    apply_komm_one_theme(prs)

    # Save
    os.makedirs(output_dir, exist_ok=True)
    filepath = os.path.join(output_dir, filename)
    prs.save(filepath)
    return filepath
